"""Resource-contained document text extractor. Invoked only as a child process."""
from __future__ import annotations

import itertools
import json
from pathlib import Path
import sys


MAX_CHARS = 12_000


def _apply_resource_limits() -> bool:
    """Apply limits inside this fresh single-threaded worker, before parsing."""
    try:
        import resource
    except ImportError:
        return False
    memory = 512 * 1024 * 1024
    limits = (
        ("RLIMIT_CORE", 0, 0),
        ("RLIMIT_AS", memory, memory),
        ("RLIMIT_CPU", 15, 16),
        ("RLIMIT_FSIZE", 4 * 1024 * 1024, 4 * 1024 * 1024),
        ("RLIMIT_NOFILE", 64, 64),
        ("RLIMIT_NPROC", 32, 32),
    )
    try:
        for name, soft, hard in limits:
            resource_key = getattr(resource, name, None)
            if resource_key is not None:
                resource.setrlimit(resource_key, (soft, hard))
        return True
    except (OSError, ValueError):
        return False


def _append(parts: list[str], value) -> bool:
    text = str(value or "").strip()
    if text:
        remaining = MAX_CHARS - sum(len(item) + 1 for item in parts)
        if remaining <= 0:
            return False
        parts.append(text[:remaining])
    return sum(len(item) + 1 for item in parts) < MAX_CHARS


def extract(path: Path, ext: str) -> str:
    parts: list[str] = []
    if ext == ".pdf":
        from pypdf import PdfReader

        reader = PdfReader(str(path), strict=True)
        if len(reader.pages) > 200:
            raise ValueError("PDF 页数超过限制")
        for page in reader.pages[:30]:
            if not _append(parts, page.extract_text() or ""):
                break
    elif ext == ".docx":
        from docx import Document

        document = Document(str(path))
        for paragraph in document.paragraphs[:5000]:
            if not _append(parts, paragraph.text):
                break
        for table in document.tables[:200]:
            for row in table.rows[:1000]:
                if not _append(
                    parts,
                    " | ".join(
                        cell.text.strip()
                        for cell in row.cells
                        if cell.text.strip()
                    ),
                ):
                    break
    elif ext == ".xlsx":
        from openpyxl import load_workbook

        workbook = load_workbook(
            str(path), read_only=True, data_only=True, keep_links=False
        )
        if len(workbook.worksheets) > 50:
            raise ValueError("工作表数量超过限制")
        rows_seen = 0
        for worksheet in workbook.worksheets:
            _append(parts, f"【工作表:{worksheet.title}】")
            for row in worksheet.iter_rows(values_only=True):
                rows_seen += 1
                if rows_seen > 5000:
                    raise ValueError("表格行数超过限制")
                if not _append(
                    parts,
                    " | ".join(
                        str(cell).strip()
                        for cell in row
                        if cell is not None and str(cell).strip()
                    ),
                ):
                    break
    elif ext == ".pptx":
        from pptx import Presentation

        presentation = Presentation(str(path))
        if len(presentation.slides) > 200:
            raise ValueError("幻灯片页数超过限制")
        # python-pptx 的 Slides/Shapes 不支持切片(会抛 AttributeError),必须用
        # islice 逐项取,否则任何 .pptx 解析都 100% 失败。
        for number, slide in enumerate(
            itertools.islice(presentation.slides, 100), 1
        ):
            _append(parts, f"【第{number}页】")
            for shape in itertools.islice(slide.shapes, 1000):
                if (
                    getattr(shape, "has_text_frame", False)
                    and not _append(parts, shape.text_frame.text)
                ):
                    break
    else:
        raise ValueError("unsupported extension")
    return "\n".join(parts)[:MAX_CHARS]


def main(argv: list[str]) -> int:
    if len(argv) != 3 or argv[2] not in {".pdf", ".docx", ".xlsx", ".pptx"}:
        return 2
    if not _apply_resource_limits():
        return 2
    try:
        text = extract(Path(argv[1]), argv[2])
        sys.stdout.write(json.dumps({"text": text}, ensure_ascii=False))
        return 0
    except Exception:
        return 1


if __name__ == "__main__":
    raise SystemExit(main(sys.argv))
